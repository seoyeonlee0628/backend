from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG=RGBColor(0x0f,0x17,0x2a); CYAN=RGBColor(0x06,0xb6,0xd4); PURPLE=RGBColor(0x8b,0x5c,0xf6)
WHITE=RGBColor(0xff,0xff,0xff); GRAY=RGBColor(0x94,0xa3,0xb8); DARK2=RGBColor(0x1e,0x29,0x3b)
GREEN=RGBColor(0x22,0xc5,0x5e); YELLOW=RGBColor(0xf5,0x9e,0x0b)
W=Cm(33.87); H=Cm(19.05)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
blank=prs.slide_layouts[6]

def sl_new():
    s=prs.slides.add_slide(blank); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=BG; return s

def rc(s,x,y,w,h,c):
    sh=s.shapes.add_shape(1,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background(); return sh

def tx(s,x,y,w,h,t,sz,bold=False,color=WHITE,align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=t; r.font.name="맑은 고딕"; r.font.size=Pt(sz); r.font.bold=bold
    r.font.color.rgb=color; r.font.italic=italic; return tb

def bul(s,x,y,w,h,items,sz=12,color=GRAY):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; first=True
    for item in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text="• "+item
        r.font.name="맑은 고딕"; r.font.size=Pt(sz); r.font.color.rgb=color; p.space_before=Pt(3)

def base(s,num):
    rc(s,Cm(0),Cm(0),Cm(0.35),H,CYAN); rc(s,Cm(0),H-Cm(0.25),W,Cm(0.25),CYAN)
    tx(s,W-Cm(2.5),H-Cm(0.95),Cm(2.2),Cm(0.65),f"{num}/12",9,color=GRAY,align=PP_ALIGN.RIGHT)

def hdr(s,t1,t2):
    tx(s,Cm(1.2),Cm(0.8),Cm(28),Cm(1.2),t1,22,bold=True,color=CYAN)
    tx(s,Cm(1.2),Cm(2.1),Cm(28),Cm(0.75),t2,11,color=PURPLE,italic=True)
    rc(s,Cm(1.2),Cm(2.95),Cm(12),Cm(0.05),PURPLE)

def card(s,x,y,w,h,title,body,color=CYAN):
    sh=s.shapes.add_shape(1,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=DARK2
    sh.line.color.rgb=color; sh.line.width=Pt(1); rc(s,x,y,w,Cm(0.2),color)
    tx(s,x+Cm(0.3),y+Cm(0.35),w-Cm(0.6),Cm(0.7),title,11,bold=True,color=color)
    tx(s,x+Cm(0.3),y+Cm(1.1),w-Cm(0.6),h-Cm(1.2),body,10.5,color=GRAY)

# S1 Title
s=sl_new()
rc(s,Cm(0),Cm(0),Cm(0.5),H,CYAN); rc(s,Cm(0),H-Cm(0.35),W,Cm(0.35),CYAN)
sh=s.shapes.add_shape(1,Cm(5),Cm(4.2),Cm(24),Cm(10.8))
sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(0x06,0x2a,0x35); sh.line.color.rgb=CYAN; sh.line.width=Pt(1)
tx(s,Cm(6),Cm(5.0),Cm(22),Cm(2.5),"복지체크",54,bold=True,color=CYAN,align=PP_ALIGN.CENTER)
tx(s,Cm(6),Cm(7.7),Cm(22),Cm(1.1),"AI 기반 청년 복지 통합 플랫폼",21,color=WHITE,align=PP_ALIGN.CENTER)
tx(s,Cm(6),Cm(9.0),Cm(22),Cm(0.9),"취준생을 위한 맞춤형 복지 자가진단 서비스",13,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
tags=["Spring Boot 3.x","Python FastAPI","MySQL 8.0","Thymeleaf SPA"]; qx=Cm(7)
for tag in tags:
    sh2=s.shapes.add_shape(1,qx,Cm(12.2),Cm(4.4),Cm(0.72))
    sh2.fill.solid(); sh2.fill.fore_color.rgb=RGBColor(0x06,0x2a,0x35); sh2.line.color.rgb=CYAN; sh2.line.width=Pt(0.75)
    tx(s,qx+Cm(0.1),Cm(12.2),Cm(4.2),Cm(0.72),tag,9,color=CYAN,align=PP_ALIGN.CENTER); qx+=Cm(5.0)
tx(s,Cm(1),H-Cm(0.95),Cm(22),Cm(0.6),"2026.06  |  취준생 전용 복지 플랫폼 프로젝트",9,color=GRAY)

# S2 프로젝트 소개
s=sl_new(); base(s,2); hdr(s,"프로젝트 소개","복지체크 서비스 개요")
tx(s,Cm(1.2),Cm(3.4),Cm(30),Cm(0.8),"복지체크는 취준생을 위한 AI 기반 원스톱 복지 자가진단 플랫폼입니다.",14,bold=True,color=WHITE)
cards=[("서비스 목적","복잡한 복지 정보를 쉽게 찾고\n신청할 수 있도록 원스톱 제공"),("대상 사용자","취준생 (19~34세)\n청년 구직자 / 사회 초년생"),("핵심 가치","AI 맞춤 추천 + 접근성 향상\n지속적 동기부여 게이미피케이션"),("플랫폼 특징","AI 진단+로드맵+커뮤니티\n+찜 캘린더 완전 통합")]
cx=Cm(1.2)
for ct,cb in cards: card(s,cx,Cm(4.5),Cm(7.5),Cm(6.5),ct,cb); cx+=Cm(8.0)
bul(s,Cm(1.2),Cm(11.8),Cm(30),Cm(5.5),["자연어 입력 한 줄 → AI 맞춤 복지 서비스 즉시 추천","7개 직군 × 3기간 취업 로드맵 자동 생성","출석 포인트 → 뱃지·이모지·특별 직군 해금 게이미피케이션","커뮤니티, 찜 캘린더, 마이페이지 통합 제공"],12.5,WHITE)

# S3 개발 배경
s=sl_new(); base(s,3); hdr(s,"개발 배경","왜 이 서비스가 필요한가")
stats=[("42%","청년 복지\n수혜율"),("2.3만개+","복지 서비스\n총 개수"),("73%","정보 탐색\n어렵다 응답"),("68%","복지 신청\n포기 경험")]
sx=Cm(1.2)
for val,lbl in stats:
    sh=s.shapes.add_shape(1,sx,Cm(3.5),Cm(7.0),Cm(4.8))
    sh.fill.solid(); sh.fill.fore_color.rgb=DARK2; sh.line.color.rgb=CYAN; sh.line.width=Pt(1)
    tx(s,sx,Cm(3.7),Cm(7.0),Cm(2.0),val,28,bold=True,color=CYAN,align=PP_ALIGN.CENTER)
    tx(s,sx,Cm(5.7),Cm(7.0),Cm(1.5),lbl,11,color=GRAY,align=PP_ALIGN.CENTER); sx+=Cm(7.8)
tx(s,Cm(1.2),Cm(9.0),Cm(14),Cm(0.65),"핵심 문제점",13,bold=True,color=WHITE)
bul(s,Cm(1.2),Cm(9.8),Cm(14),Cm(5.5),["정보 분산: 복지로·워크넷·지자체 사이트 각각","복잡한 신청 절차로 중도 포기율 높음","개인 상황 맞춤 추천 서비스 전무","취업 준비와 복지 정보 연계 부재"])
tx(s,Cm(16.5),Cm(9.0),Cm(14),Cm(0.65),"복지체크의 해결책",13,bold=True,color=CYAN)
bul(s,Cm(16.5),Cm(9.8),Cm(14),Cm(5.5),["자연어 1줄 → AI 복지 즉시 추천","취업 로드맵 + 복지 정보 원스톱 통합","출석 포인트로 지속 재방문 유도","모바일 친화적 SPA로 접근성 극대화"],12,GRAY)

# S4 시장 조사
s=sl_new(); base(s,4); hdr(s,"시장 조사","기존 서비스 현황 분석")
tx(s,Cm(1.2),Cm(3.5),Cm(30),Cm(0.65),"기존 청년 복지·취업 플랫폼 비교 분석",13,bold=True,color=WHITE)
cols2=[Cm(1.2),Cm(8.5),Cm(17.0),Cm(25.2)]; wids=[Cm(7.0),Cm(8.2),Cm(7.9),Cm(7.0)]
hdrs2=["서비스","주요 기능","한계점","대상"]
for cx2,cw2,hd2 in zip(cols2,wids,hdrs2):
    rc(s,cx2,Cm(4.4),cw2,Cm(0.72),CYAN); tx(s,cx2+Cm(0.1),Cm(4.4),cw2-Cm(0.2),Cm(0.72),hd2,10,bold=True,color=BG,align=PP_ALIGN.CENTER)
rows2=[("복지로","복지 정보 통합 제공\n온라인 신청 지원","개인화 없음\nAI 추천 전무","전 국민"),("워크넷","구인·구직 매칭\n취업 정보 제공","복지 연계 없음\n청년 특화 부족","구직자"),("청년정책","청년 정책 정보\n지자체 연계","정적 정보 나열\n상호작용 없음","청년")]
ry2=Cm(5.2)
for row2 in rows2:
    for cx2,cw2,rd2 in zip(cols2,wids,row2):
        rc(s,cx2,ry2,cw2,Cm(1.8),DARK2); tx(s,cx2+Cm(0.15),ry2+Cm(0.1),cw2-Cm(0.3),Cm(1.6),rd2,10,color=GRAY)
    ry2+=Cm(1.9)
hl=["복지체크","AI 진단+로드맵\n커뮤니티+캘린더","없음 (우리가 해결)","취준생 특화"]; hc=[CYAN,WHITE,GREEN,CYAN]
for cx2,cw2,rd2,hcc in zip(cols2,wids,hl,hc):
    sh=s.shapes.add_shape(1,cx2,ry2,cw2,Cm(1.8)); sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(0x06,0x2a,0x35); sh.line.color.rgb=CYAN; sh.line.width=Pt(1.2)
    tx(s,cx2+Cm(0.15),ry2+Cm(0.1),cw2-Cm(0.3),Cm(1.6),rd2,10,bold=True,color=hcc)
tx(s,Cm(1.2),Cm(14.5),Cm(30),Cm(0.65),"→ 기존 서비스는 복지 OR 취업 중 하나만 제공. 복지체크는 AI로 둘을 통합한 유일한 플랫폼.",11,color=YELLOW,italic=True)

# S5 차별화 전략
s=sl_new(); base(s,5); hdr(s,"차별화 전략","복지체크만의 경쟁 우위")
diffs=[("AI 자연어 진단","Python FastAPI + NLP\n자연어 1줄 → 맞춤 복지 즉시 추천\n키워드 추출 → 조건 매칭",CYAN),("직군별 로드맵","7개 직군 x 3기간 = 21가지 조합\n개발자/디자이너/마케터/기획자/공무원\n프리랜서/유튜버 포인트 해금",PURPLE),("게이미피케이션","매일 출석 +10P 포인트 적립\n뱃지·이모지·직군 해금 보상\n지속 재방문 동기 부여",CYAN),("커뮤니티 통합","취준생 전용 게시판\n좋아요·북마크·댓글·부스트\n정보 공유 + 네트워킹",PURPLE),("찜 캘린더","혜택 마감일 캘린더 관리\n찜한 혜택 한눈에 확인\n놓치는 혜택 Zero",CYAN),("보안·안정성","Spring Security 세션 인증\nCSRF 보호, 권한 접근 제어\nAPI 29개 전수 테스트 통과",PURPLE)]
for i,(dt,db,dc) in enumerate(diffs):
    card(s,Cm(1.2)+(i%3)*Cm(10.7),Cm(3.5)+(i//3)*Cm(5.6),Cm(10.2),Cm(5.1),dt,db,dc)

# S6 시스템 아키텍처
s=sl_new(); base(s,6); hdr(s,"시스템 아키텍처","전체 시스템 구조 및 데이터 흐름")
arch=[(Cm(1.2),Cm(3.7),Cm(7.8),Cm(11.5),"FRONTEND","Thymeleaf SPA\nTailwind CSS\napp.js / api.js\napi-bridge.js",CYAN),(Cm(12.0),Cm(3.7),Cm(7.8),Cm(5.3),"BACKEND","Spring Boot 3.x\nJava 17 / Port 8080\nSpring Security\nJPA / Hibernate",PURPLE),(Cm(12.0),Cm(9.5),Cm(7.8),Cm(5.3),"ML SERVER","Python FastAPI\nPort 8000\nNLP 파이프라인\n복지 키워드 매칭",YELLOW),(Cm(22.5),Cm(3.7),Cm(9.5),Cm(11.5),"DATABASE","MySQL 8.0\nPort 3306\nUser / Post / Comment\nPostLike / Bookmark\nWishItem / Attendance",GREEN)]
for bx,by,bw,bh,at,ab,ac in arch:
    sh=s.shapes.add_shape(1,bx,by,bw,bh); sh.fill.solid(); sh.fill.fore_color.rgb=DARK2; sh.line.color.rgb=ac; sh.line.width=Pt(1.5)
    rc(s,bx,by,bw,Cm(0.25),ac); tx(s,bx+Cm(0.3),by+Cm(0.4),bw-Cm(0.6),Cm(0.8),at,12,bold=True,color=ac)
    tx(s,bx+Cm(0.3),by+Cm(1.3),bw-Cm(0.6),bh-Cm(1.5),ab,10.5,color=GRAY)
for ax,ay,at2,ac2 in [(Cm(9.3),Cm(6.0),"HTTP <-> JSON",CYAN),(Cm(9.3),Cm(10.5),"REST <-> JSON",YELLOW),(Cm(20.2),Cm(6.0),"JPA <-> SQL",GREEN),(Cm(20.2),Cm(10.5),"JPA <-> SQL",GREEN)]:
    tx(s,ax,ay,Cm(2.5),Cm(0.75),at2,9,bold=True,color=ac2,align=PP_ALIGN.CENTER)
tx(s,Cm(1.2),Cm(16.0),Cm(30),Cm(0.6),"사용자 -> Spring Boot API -> ML 서버 NLP 분석 -> MySQL 조회 -> JSON 응답 -> Thymeleaf 렌더링",10,color=GRAY,italic=True)

# S7 기술 구성
s=sl_new(); base(s,7); hdr(s,"기술 구성","사용된 기술 스택 상세")
groups=[("백엔드",[("Spring Boot 3.x","Java 17 기반 메인 API 서버"),("Spring Security","세션 인증, CSRF 보호, 권한 관리"),("JPA / Hibernate","ORM, MySQL 연동"),("Maven","빌드 도구 및 의존성 관리")],PURPLE),("프론트엔드",[("Thymeleaf","서버사이드 템플릿 엔진"),("Tailwind CSS","유틸리티 CSS 프레임워크"),("Vanilla JS (SPA)","goPage() 기반 단일 페이지"),("api-bridge.js","서버 API 오버라이드 레이어")],CYAN),("AI / ML",[("Python FastAPI","ML 추론 API 서버 (Port 8000)"),("NLP 파이프라인","키워드 추출 + 복지 조건 매칭"),("roadmap_templates","21가지 로드맵 템플릿 관리"),("LocalModelService","Spring <-> FastAPI 연동")],YELLOW),("인프라 / DB",[("MySQL 8.0","메인 관계형 DB (Port 3306)"),("JPA 엔티티","User/Post/Comment/Wish 등"),("Spring DevTools","개발 환경 핫리로드"),("REST API 29개","전수 자동화 테스트 통과")],GREEN)]
for gi,(gn,gi_items,gc) in enumerate(groups):
    gx=Cm(1.2)+(gi%2)*Cm(16.0); gy=Cm(3.5)+(gi//2)*Cm(7.2)
    tx(s,gx,gy,Cm(15),Cm(0.72),gn,13,bold=True,color=gc); rc(s,gx,gy+Cm(0.72),Cm(14.5),Cm(0.05),gc)
    for j,(tech,desc) in enumerate(gi_items):
        iy=gy+Cm(0.95)+j*Cm(1.35); rc(s,gx,iy+Cm(0.2),Cm(0.22),Cm(0.5),gc)
        tx(s,gx+Cm(0.4),iy,Cm(5.0),Cm(0.7),tech,11,bold=True,color=WHITE); tx(s,gx+Cm(5.6),iy,Cm(9.5),Cm(0.7),desc,10.5,color=GRAY)

# S8 핵심 기능 구현
s=sl_new(); base(s,8); hdr(s,"핵심 기능 구현","맞춤형 복지 추천 및 취업 지원 기능")
feats=[("AI 복지 진단","자연어 입력 -> Python NLP\n맞춤 복지 서비스 추천\n키워드 추출+조건 매칭",CYAN),("취업 로드맵","7개 직군 x 3기간 = 21가지\n월별 세부 목표 자동 생성\n프리랜서/유튜버 80P 해금",PURPLE),("출석 포인트","매일 출석 +10P 자동 지급\n이중 출석 방지 로직\n연속 출석 스트릭 표시",CYAN),("뱃지 시스템","칭호 뱃지 / 특별 직군 해금\n닉네임 이모지 (20P)\n뱃지 상점 실시간 반영",PURPLE),("커뮤니티","게시글 CRUD+좋아요/북마크\n댓글+게시글 부스트(50P)\nauthorEmoji 실시간 표시",CYAN),("찜 캘린더","복지 혜택 찜+마감일 관리\n캘린더 UI 한눈 확인\n찜 목록 마이페이지 통합",PURPLE)]
for i,(ft,fb,fc) in enumerate(feats):
    card(s,Cm(1.2)+(i%3)*Cm(10.7),Cm(3.5)+(i//3)*Cm(5.6),Cm(10.2),Cm(5.1),ft,fb,fc)

# S9 서비스 구성
s=sl_new(); base(s,9); hdr(s,"서비스 구성","페이지별 화면 구성")
pages=[("홈","서비스 소개\n출석 체크\n빠른 메뉴"),("복지 진단","자연어 입력\nAI 분석 결과\n복지 카드"),("로드맵","직군 태그 선택\n기간 선택\n월별 목표"),("커뮤니티","게시판 목록\n좋아요/댓글\n북마크"),("찜 캘린더","혜택 찜 목록\n마감일 관리\n캘린더 UI"),("마이페이지","프로필/포인트\n뱃지/이모지\n활동 내역"),("관리자","사용자 관리\n게시글 관리\n공지 PIN")]
ppx=Cm(1.2)
for i,(pn,pd) in enumerate(pages):
    pc=CYAN if i%2==0 else PURPLE
    sh=s.shapes.add_shape(1,ppx,Cm(3.5),Cm(4.2),Cm(4.0)); sh.fill.solid(); sh.fill.fore_color.rgb=DARK2; sh.line.color.rgb=pc; sh.line.width=Pt(0.75)
    tx(s,ppx+Cm(0.1),Cm(3.7),Cm(4.0),Cm(0.75),pn,10.5,bold=True,color=pc,align=PP_ALIGN.CENTER)
    tx(s,ppx+Cm(0.1),Cm(4.6),Cm(4.0),Cm(2.7),pd,9.5,color=GRAY,align=PP_ALIGN.CENTER); ppx+=Cm(4.5)
tx(s,Cm(1.2),Cm(8.3),Cm(30),Cm(0.65),"JavaScript 레이어 구조",13,bold=True,color=WHITE)
for jn,jd,jc in [("api.js","서버 HTTP 통신 함수 (apiFetch, apiLogin, apiCreatePost 등)",CYAN),("app.js","UI 렌더링 및 이벤트 핸들러 (toggleJobTag, buyEmoji, openMyPage 등)",WHITE),("api-bridge.js","서버 연동 오버라이드 레이어 (submitPost, toggleLikePost, openPost 등)",PURPLE)]:
    jy={"api.js":Cm(9.1),"app.js":Cm(9.95),"api-bridge.js":Cm(10.8)}[jn]
    rc(s,Cm(1.2),jy,Cm(3.5),Cm(0.65),DARK2); tx(s,Cm(1.3),jy,Cm(3.3),Cm(0.65),jn,10,bold=True,color=jc,align=PP_ALIGN.CENTER)
    tx(s,Cm(5.0),jy,Cm(26),Cm(0.65),"->  "+jd,10,color=GRAY)
tx(s,Cm(1.2),Cm(11.8),Cm(30),Cm(0.65),"DB 엔티티",13,bold=True,color=WHITE)
tx(s,Cm(1.2),Cm(12.5),Cm(30),Cm(0.65),"User | Post | Comment | PostLike | Bookmark | WishItem | Attendance | Board",11,color=GRAY)

# S10 개발 과정
s=sl_new(); base(s,10); hdr(s,"개발 과정 (마일스톤)","단계별 개발 프로세스")
ms=[("M1","요구사항 분석\n& DB 설계","엔티티 설계\n관계 정의\nERD 작성",CYAN),("M2","Spring Boot\nAPI 서버 구축","인증/인가\nREST API\nJPA 연동",PURPLE),("M3","Python ML\n서버 연동","FastAPI 구축\nNLP 파이프라인\nSpring 연동",YELLOW),("M4","Thymeleaf SPA\n프론트 구현","goPage SPA\nUI 컴포넌트\napi.js 레이어",CYAN),("M5","게이미피케이션\n& 특별 기능","포인트/뱃지\n직군 해금\n이모지 시스템",PURPLE),("M6","통합 테스트\n& 완성","29/29 PASS\n버그 수정\n최종 완성",GREEN)]
mmx=Cm(1.2)
for ii,(mn,mt,md,mc) in enumerate(ms):
    if ii<5: tx(s,mmx+Cm(5.1),Cm(6.5),Cm(0.6),Cm(0.7),"->",13,bold=True,color=GRAY,align=PP_ALIGN.CENTER)
    sh=s.shapes.add_shape(1,mmx,Cm(3.8),Cm(4.9),Cm(8.5)); sh.fill.solid(); sh.fill.fore_color.rgb=DARK2; sh.line.color.rgb=mc; sh.line.width=Pt(1.2)
    rc(s,mmx,Cm(3.8),Cm(4.9),Cm(0.22),mc)
    tx(s,mmx+Cm(0.1),Cm(4.1),Cm(4.7),Cm(0.72),mn,13,bold=True,color=mc,align=PP_ALIGN.CENTER)
    tx(s,mmx+Cm(0.1),Cm(4.9),Cm(4.7),Cm(1.5),mt,10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tx(s,mmx+Cm(0.1),Cm(6.5),Cm(4.7),Cm(3.0),md,9.5,color=GRAY,align=PP_ALIGN.CENTER); mmx+=Cm(5.4)
tx(s,Cm(1.2),Cm(13.3),Cm(30),Cm(0.65),"주요 기술 이슈 해결",13,bold=True,color=WHITE)
bul(s,Cm(1.2),Cm(14.1),Cm(30),Cm(3.5),["api-bridge.js 오버라이드 패턴 -> app.js 서버 미연동 문제 해결","PostService toggleLike() Map<> 반환 -> likeCount 실시간 반영","PowerShell $PID 예약 변수 충돌 -> $postIdVal 로 변경, 29/29 전수 통과","Python NLP 7개 직군 로드맵 템플릿 완성 (프리랜서/유튜버 신규 추가)"],11,GRAY)

# S11 성과
s=sl_new(); base(s,11); hdr(s,"성과","개발 완료 지표 및 달성 현황")
mets=[("29 / 29","REST API 전수\n테스트 통과",GREEN),("7개","직군별\n취업 로드맵",CYAN),("21가지","로드맵 조합\n(직군x기간)",PURPLE),("6단계","마일스톤\n전체 완료",YELLOW)]
mmx2=Cm(1.2)
for mv,ml,mc in mets:
    sh=s.shapes.add_shape(1,mmx2,Cm(3.5),Cm(7.5),Cm(4.8)); sh.fill.solid(); sh.fill.fore_color.rgb=DARK2; sh.line.color.rgb=mc; sh.line.width=Pt(1.2)
    tx(s,mmx2,Cm(3.7),Cm(7.5),Cm(2.2),mv,28,bold=True,color=mc,align=PP_ALIGN.CENTER)
    tx(s,mmx2,Cm(5.8),Cm(7.5),Cm(1.5),ml,11,color=GRAY,align=PP_ALIGN.CENTER); mmx2+=Cm(8.0)
tx(s,Cm(1.2),Cm(9.3),Cm(30),Cm(0.65),"구현 완료 기능 목록",13,bold=True,color=WHITE)
bul(s,Cm(1.2),Cm(10.1),Cm(30),Cm(6.5),["AI 복지 진단: 자연어 입력 -> Python NLP -> 맞춤 복지 서비스 추천","취업 로드맵: 7개 직군 x 3기간 월별 세부 목표 자동 생성","포인트 시스템: 출석 +10P, 뱃지/이모지 구매, 특별 직군 80P 해금","커뮤니티: 게시글 CRUD, 좋아요/북마크/댓글/부스트, 닉네임 이모지 표시","찜 캘린더: 복지 혜택 찜 + 마감일 캘린더 관리","관리자: 사용자 관리, 게시글 관리, 공지 PIN 설정"],12,GRAY)

# S12 향후 발전 방향
s=sl_new(); base(s,12); hdr(s,"향후 발전 방향","서비스 고도화 및 확장 계획 수립")
plans=[("단기 (3개월)",["정부 복지로 Open API 실연동","지역별/나이별 복지 필터링 강화","소셜 로그인 (카카오/네이버)","PWA 변환으로 모바일 앱 경험"],CYAN),("중기 (6개월)",["GPT 기반 대화형 복지 상담 챗봇","로드맵 AI 개인화 (진단 이력 학습)","기업 취업 공고 API 연동","복지 신청 현황 트래킹 기능"],PURPLE),("장기 (1년+)",["React Native 네이티브 앱 출시","복지 수급자 멘토링 매칭 시스템","B2G 지자체 복지 플랫폼 납품","ML 고도화: 성공 사례 데이터 학습"],YELLOW)]
prx=Cm(1.2)
for pp,pi,pc in plans:
    sh=s.shapes.add_shape(1,prx,Cm(3.5),Cm(10.0),Cm(11.2)); sh.fill.solid(); sh.fill.fore_color.rgb=DARK2; sh.line.color.rgb=pc; sh.line.width=Pt(1.2)
    rc(s,prx,Cm(3.5),Cm(10.0),Cm(0.28),pc)
    tx(s,prx+Cm(0.3),Cm(3.9),Cm(9.5),Cm(0.75),pp,13,bold=True,color=pc)
    bul(s,prx+Cm(0.2),Cm(4.8),Cm(9.5),Cm(9.5),pi,11.5,GRAY); prx+=Cm(10.8)
sh=s.shapes.add_shape(1,Cm(1.2),Cm(15.5),Cm(31.5),Cm(2.0)); sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(0x06,0x2a,0x35); sh.line.color.rgb=CYAN; sh.line.width=Pt(1)
tx(s,Cm(1.5),Cm(15.75),Cm(31),Cm(1.5),"복지체크는 취준생의 복지 접근성을 높이고, AI와 게이미피케이션으로 지속 사용을 유도하는 플랫폼으로 발전시켜 나가겠습니다.",11.5,color=WHITE,align=PP_ALIGN.CENTER,italic=True)

out=r"C:\Users\dd\Downloads\welfare-check\복지체크_발표자료.pptx"
prs.save(out)
print(f"OK: {out}")
